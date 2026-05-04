"""
PPTX Slide Parser — extracts all content elements from a single PowerPoint slide.

Uses python-pptx to read shapes, text, images, tables, and their raw XML.
Produces a structured ``ParsedSlide`` dataclass that downstream modules consume.

Corresponds to **Step 1-2** of the workflow: "Users upload a PPTX slide → AI parses content".
"""

from __future__ import annotations

import io
import re
import base64
from pathlib import Path
from dataclasses import dataclass, field
from typing import Any
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ---------------------------------------------------------------------------
# Data classes — the parsed representation of a slide
# ---------------------------------------------------------------------------

@dataclass
class TextRun:
    """A single formatted run of text within a paragraph."""
    text: str
    font_name: str | None = None
    font_size_pt: float | None = None
    bold: bool = False
    italic: bool = False
    underline: bool = False
    color_hex: str | None = None  # "RRGGBB" without #
    fill_opacity: float = 1.0


@dataclass
class ParagraphInfo:
    """A paragraph (one or more TextRuns) with alignment."""
    runs: list[TextRun] = field(default_factory=list)
    alignment: str | None = None  # LEFT, CENTER, RIGHT, JUSTIFY
    line_spacing: float | None = None
    space_before_pt: float = 0.0
    space_after_pt: float = 0.0
    bullet: bool = False
    level: int = 0

    @property
    def plain_text(self) -> str:
        return "".join(r.text for r in self.runs)


@dataclass
class ShapeInfo:
    """Any shape on the slide — text box, rectangle, image, group, table, chart."""
    shape_id: int
    name: str
    shape_type: str  # "text_box", "rectangle", "image", "table", "group", "chart", "line", etc.
    left_emu: int
    top_emu: int
    width_emu: int
    height_emu: int

    # Text content (only for text-bearing shapes)
    paragraphs: list[ParagraphInfo] = field(default_factory=list)

    # Fill
    fill_color_hex: str | None = None
    fill_opacity: float = 1.0

    # Stroke
    stroke_color_hex: str | None = None
    stroke_width_emu: int = 0

    # Rounding
    corner_radius_emu: int = 0  # for rounded rectangles

    # Image
    image_blob: bytes | None = None
    image_ext: str | None = None  # "png", "jpeg", etc.

    # Table data (list of rows, each row is list of cell strings)
    table_data: list[list[str]] = field(default_factory=list)

    # Raw XML for deep inspection
    raw_xml: str = ""


@dataclass
class ParsedSlide:
    """Complete parsed representation of a single PowerPoint slide."""
    slide_index: int = 0
    slide_width_emu: int = 0
    slide_height_emu: int = 0
    slide_width_px: float = 0.0
    slide_height_px: float = 0.0
    shapes: list[ShapeInfo] = field(default_factory=list)
    notes: str = ""

    # Convenience accessors
    @property
    def title_shape(self) -> ShapeInfo | None:
        """Heuristic: the topmost text-bearing shape is the title."""
        text_shapes = sorted(
            [s for s in self.shapes if s.paragraphs],
            key=lambda s: s.top_emu,
        )
        return text_shapes[0] if text_shapes else None

    @property
    def text_shapes(self) -> list[ShapeInfo]:
        return [s for s in self.shapes if s.paragraphs]

    @property
    def image_shapes(self) -> list[ShapeInfo]:
        return [s for s in self.shapes if s.shape_type == "image"]


# ---------------------------------------------------------------------------
# EMU ↔ px conversion (96 DPI, matching ppt-master convention)
# ---------------------------------------------------------------------------

EMU_PER_PX = 9525

def emu_to_px(emu: int) -> float:
    return emu / EMU_PER_PX

def px_to_emu(px: float) -> int:
    return int(round(px * EMU_PER_PX))


# ---------------------------------------------------------------------------
# Color extraction helpers
# ---------------------------------------------------------------------------

def _extract_hex_from_element(el: ET.Element | None) -> str | None:
    """Extract #RRGGBB or RRGGBB from a DrawingML element like <a:solidFill>."""
    if el is None:
        return None
    xml_str = ET.tostring(el, encoding="unicode")
    m = re.search(r'(?:val|lastClr)="([0-9A-Fa-f]{6})"', xml_str)
    if m:
        return m.group(1).upper()
    return None


def _extract_fill_info(shape_xml: ET.Element) -> tuple[str | None, float]:
    """Return (hex_color, opacity) from a shape's <a:solidFill> or <a:gradFill>."""
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    solid = shape_xml.find(".//a:solidFill/a:srgbClr", ns)
    if solid is not None:
        color = solid.get("val", "").upper() or None
        alpha_el = solid.find("a:alpha", ns)
        opacity = int(alpha_el.get("val", "100000")) / 100000.0 if alpha_el is not None else 1.0
        return color, opacity
    # For gradients, take the first stop color
    grad = shape_xml.find(".//a:gradFill", ns)
    if grad is not None:
        first_stop = grad.find(".//a:gs/a:srgbClr", ns)
        if first_stop is not None:
            return first_stop.get("val", "").upper() or None, 1.0
    return None, 1.0


# ---------------------------------------------------------------------------
# Main parser
# ---------------------------------------------------------------------------

def parse_pptx_slide(pptx_path: str | Path, slide_index: int = 0) -> ParsedSlide:
    """Parse a single slide from a PPTX file into a ``ParsedSlide``.

    Args:
        pptx_path: Path to a .pptx file.
        slide_index: 0-based index of the slide to parse.

    Returns:
        ParsedSlide with all shapes, text, images, and metadata extracted.
    """
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))

    if slide_index >= len(prs.slides):
        raise ValueError(
            f"Slide index {slide_index} out of range (file has {len(prs.slides)} slide(s))"
        )

    slide = prs.slides[slide_index]

    parsed = ParsedSlide(
        slide_index=slide_index,
        slide_width_emu=prs.slide_width,
        slide_height_emu=prs.slide_height,
        slide_width_px=emu_to_px(prs.slide_width),
        slide_height_px=emu_to_px(prs.slide_height),
    )

    # Notes
    try:
        notes_slide = slide.notes_slide
        parsed.notes = notes_slide.notes_text_frame.text if notes_slide else ""
    except Exception:
        pass

    for shape in slide.shapes:
        si = _parse_shape(shape)
        if si:
            parsed.shapes.append(si)

    return parsed


def _parse_shape(shape: Any) -> ShapeInfo | None:
    """Route a python-pptx shape to the appropriate parsing logic."""
    try:
        raw_xml = ET.tostring(shape._element, encoding="unicode")
    except Exception:
        raw_xml = ""

    fill_color, fill_opacity = _extract_fill_info(shape._element)

    # Determine shape type
    if shape.has_table:
        return _parse_table_shape(shape, raw_xml, fill_color, fill_opacity)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return _parse_group_shape(shape, raw_xml)
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return _parse_image_shape(shape, raw_xml)
    elif shape.has_text_frame:
        return _parse_text_shape(shape, raw_xml, fill_color, fill_opacity)
    else:
        return _parse_simple_shape(shape, raw_xml, fill_color, fill_opacity)


def _parse_text_shape(
    shape: Any, raw_xml: str, fill_color: str | None, fill_opacity: float
) -> ShapeInfo:
    """Parse a text-bearing shape (text box, rectangle with text, etc.)."""
    tf = shape.text_frame
    paragraphs = []
    for para in tf.paragraphs:
        runs = []
        for run in para.runs:
            color_hex = None
            try:
                run_fill = run.font.fill
                if run_fill and run_fill.type is not None:
                    color_hex = str(run_fill.fore_color.rgb) if run_fill.fore_color else None
            except Exception:
                pass
            runs.append(TextRun(
                text=run.text,
                font_name=run.font.name,
                font_size_pt=run.font.size / 12700 if run.font.size else None,
                bold=run.font.bold or False,
                italic=run.font.italic or False,
                underline=run.font.underline or False,
                color_hex=color_hex,
            ))
        paragraphs.append(ParagraphInfo(
            runs=runs,
            alignment=str(para.alignment) if para.alignment else None,
            line_spacing=para.line_spacing,
            space_before_pt=para.space_before,
            space_after_pt=para.space_after,
            bullet=para.level is not None and para.level > 0,
            level=para.level or 0,
        ))

    shape_type = "text_box"
    # Detect rounded rectangle
    try:
        prstGeom = shape._element.find(
            './/{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom'
        )
        if prstGeom is not None and prstGeom.get("prst") == "roundRect":
            shape_type = "rounded_rectangle"
    except Exception:
        pass

    return ShapeInfo(
        shape_id=shape.shape_id,
        name=shape.name,
        shape_type=shape_type,
        left_emu=shape.left,
        top_emu=shape.top,
        width_emu=shape.width,
        height_emu=shape.height,
        paragraphs=paragraphs,
        fill_color_hex=fill_color,
        fill_opacity=fill_opacity,
        raw_xml=raw_xml,
    )


def _parse_image_shape(shape: Any, raw_xml: str) -> ShapeInfo:
    """Parse an image/picture shape, extracting the raw image bytes."""
    image_blob = None
    image_ext = None
    try:
        image = shape.image
        image_blob = image.blob
        image_ext = image.content_type.split("/")[-1] if image.content_type else "png"
    except Exception:
        pass

    return ShapeInfo(
        shape_id=shape.shape_id,
        name=shape.name,
        shape_type="image",
        left_emu=shape.left,
        top_emu=shape.top,
        width_emu=shape.width,
        height_emu=shape.height,
        image_blob=image_blob,
        image_ext=image_ext,
        raw_xml=raw_xml,
    )


def _parse_table_shape(
    shape: Any, raw_xml: str, fill_color: str | None, fill_opacity: float
) -> ShapeInfo:
    """Parse a table shape."""
    table = shape.table
    table_data: list[list[str]] = []
    for row in table.rows:
        row_data: list[str] = []
        for cell in row.cells:
            row_data.append(cell.text)
        table_data.append(row_data)

    return ShapeInfo(
        shape_id=shape.shape_id,
        name=shape.name,
        shape_type="table",
        left_emu=shape.left,
        top_emu=shape.top,
        width_emu=shape.width,
        height_emu=shape.height,
        table_data=table_data,
        fill_color_hex=fill_color,
        fill_opacity=fill_opacity,
        raw_xml=raw_xml,
    )


def _parse_group_shape(shape: Any, raw_xml: str) -> ShapeInfo:
    """Parse a group shape (surface-level only; children are not recursed by default)."""
    return ShapeInfo(
        shape_id=shape.shape_id,
        name=shape.name,
        shape_type="group",
        left_emu=shape.left,
        top_emu=shape.top,
        width_emu=shape.width,
        height_emu=shape.height,
        raw_xml=raw_xml,
    )


def _parse_simple_shape(
    shape: Any, raw_xml: str, fill_color: str | None, fill_opacity: float
) -> ShapeInfo:
    """Parse a simple shape (rectangle, line, arrow, etc.) with no text."""
    shape_type = str(shape.shape_type).lower().replace("mso_shape_type.", "") if shape.shape_type else "unknown"

    # Detect line
    if shape_type in ("line", "line_inverse"):
        shape_type = "line"

    return ShapeInfo(
        shape_id=shape.shape_id,
        name=shape.name,
        shape_type=shape_type,
        left_emu=shape.left,
        top_emu=shape.top,
        width_emu=shape.width,
        height_emu=shape.height,
        fill_color_hex=fill_color,
        fill_opacity=fill_opacity,
        raw_xml=raw_xml,
    )
